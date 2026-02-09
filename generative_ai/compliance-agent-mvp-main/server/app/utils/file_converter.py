"""
Utility to convert files (PDF, DOCX, TXT, PPTX, CSV, Excel) to Markdown format.
Based on server/services/new/file_to_markdown.py
"""

import re
from pathlib import Path

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

try:
    import mammoth
except Exception:
    mammoth = None

try:
    from markdownify import markdownify as html_to_md
except Exception:
    html_to_md = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    import pandas as pd
except Exception:
    pd = None

try:
    import openpyxl
except Exception:
    openpyxl = None

try:
    import xlrd
except Exception:
    xlrd = None


def sanitize_stem(name: str) -> str:
    """Sanitize a filename stem for safe file system usage."""
    sanitized = re.sub(r"\s+", " ", name).strip()
    sanitized = sanitized.replace("/", "-")
    return sanitized


def convert_pdf_to_markdown(pdf_path: Path) -> str:
    """Convert a PDF file to Markdown format."""
    if PdfReader is None:
        raise RuntimeError(
            "pypdf is required. Please install dependencies from requirements.txt"
        )

    reader = PdfReader(str(pdf_path))
    page_texts = []
    for page in reader.pages:
        content = page.extract_text() or ""
        page_texts.append(content)
    text = "\n\n".join(page_texts)

    # Normalize line endings and whitespace
    text = text.replace("\r\n", "\n").replace("\r", "\n")

    # Basic paragraph reconstruction
    lines = [ln.rstrip() for ln in text.split("\n")]

    # Merge hyphenated line breaks
    merged_lines = []
    for i, ln in enumerate(lines):
        if ln.endswith("-") and i + 1 < len(lines):
            next_ln = lines[i + 1].lstrip()
            merged_lines.append(ln[:-1] + next_ln)
            lines[i + 1] = ""
        else:
            merged_lines.append(ln)

    # Collapse multiple blank lines
    paragraphs = []
    blank = False
    for ln in merged_lines:
        if ln.strip() == "":
            if not blank:
                paragraphs.append("")
            blank = True
        else:
            paragraphs.append(ln)
            blank = False

    body = "\n".join(paragraphs).strip()

    # Build markdown with a top-level title
    title = sanitize_stem(pdf_path.stem)
    md = f"# {title}\n\n{body}\n"
    return md


def convert_docx_to_markdown(input_path: Path) -> str:
    """Convert a DOCX file to Markdown format."""
    if mammoth is None:
        raise RuntimeError("mammoth is required. Please install dependencies from requirements.txt")
    if html_to_md is None:
        raise RuntimeError("markdownify is required. Please install dependencies from requirements.txt")

    with input_path.open("rb") as f:
        result = mammoth.convert_to_html(f)
    html = result.value or ""

    markdown = html_to_md(html, heading_style="ATX")
    if not markdown.endswith("\n"):
        markdown += "\n"
    return markdown


def convert_txt_to_markdown(input_path: Path) -> str:
    """Convert a text file to Markdown format."""
    content = input_path.read_text(encoding="utf-8", errors="ignore")
    
    # Build markdown with a top-level title
    title = sanitize_stem(input_path.stem)
    md = f"# {title}\n\n{content}\n"
    return md


def convert_pptx_to_markdown(input_path: Path) -> str:
    """Convert a PPTX file to Markdown format, preserving slide structure."""
    if Presentation is None:
        raise RuntimeError(
            "python-pptx is required. Please install dependencies from requirements.txt"
        )

    prs = Presentation(str(input_path))
    slide_texts = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        # Try to get slide title from title placeholder
        slide_title = None
        title_shape = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Check if this is a title placeholder
                if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    # Title placeholder types: TITLE (1), CENTER_TITLE (3)
                    if ph_type in (1, 3):
                        slide_title = shape.text_frame.text.strip()
                        title_shape = shape
                        break

        # Build slide header
        if slide_title:
            header = f"## Slide {slide_num}: {slide_title}"
        else:
            header = f"## Slide {slide_num}"

        # Extract all text content from shapes (skip the title shape to avoid duplication)
        content_lines = []
        for shape in slide.shapes:
            if shape is title_shape:
                continue  # Skip title shape - already included in header
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        content_lines.append(text)

        # Join content with newlines
        content = "\n".join(content_lines)
        slide_texts.append(f"{header}\n\n{content}")

    # Build markdown with a top-level title
    title = sanitize_stem(input_path.stem)
    body = "\n\n".join(slide_texts)
    md = f"# {title}\n\n{body}\n"
    return md


def _escape_markdown_table_cell(value: str) -> str:
    """Escape special characters in a markdown table cell."""
    if value is None:
        return ""
    # Convert to string and escape pipe characters
    text = str(value).replace("|", "\\|")
    # Replace newlines with spaces
    text = text.replace("\n", " ").replace("\r", " ")
    return text.strip()


def _dataframe_to_markdown_table(df, sheet_name: str = None) -> str:
    """Convert a pandas DataFrame to a markdown table."""
    if df is None or df.empty:
        return ""
    
    lines = []
    
    # Add sheet name as header if provided
    if sheet_name:
        lines.append(f"## {sheet_name}")
        lines.append("")
    
    # Get column headers
    headers = [_escape_markdown_table_cell(str(col)) for col in df.columns]
    
    # Build header row
    header_row = "| " + " | ".join(headers) + " |"
    lines.append(header_row)
    
    # Build separator row
    separator_row = "| " + " | ".join(["---"] * len(headers)) + " |"
    lines.append(separator_row)
    
    # Build data rows
    for _, row in df.iterrows():
        cells = [_escape_markdown_table_cell(str(val)) for val in row]
        data_row = "| " + " | ".join(cells) + " |"
        lines.append(data_row)
    
    return "\n".join(lines)


def convert_spreadsheet_to_markdown(input_path: Path) -> str:
    """
    Convert a spreadsheet file (CSV, XLSX, XLS, XLSM) to Markdown format.
    
    Uses pandas for consistent handling of all spreadsheet formats.
    
    Args:
        input_path: Path to the spreadsheet file
        
    Returns:
        Markdown content as a string
        
    Raises:
        RuntimeError: If required dependencies are missing or there's an error reading the file
        ValueError: If the file is empty or has no readable data
    """
    if pd is None:
        raise RuntimeError(
            "pandas is required for spreadsheet file support. "
            "Please install dependencies from requirements.txt"
        )
    
    suffix_lower = input_path.suffix.lower()
    title = sanitize_stem(input_path.stem)
    
    try:
        if suffix_lower == ".csv":
            # Read CSV file with pandas
            # Try different encodings
            encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1252']
            df = None
            last_error = None
            
            for encoding in encodings:
                try:
                    df = pd.read_csv(
                        str(input_path),
                        encoding=encoding,
                        dtype=str,  # Read all values as strings to preserve formatting
                        na_filter=False,  # Don't convert empty cells to NaN
                        on_bad_lines='warn'  # Warn but don't fail on bad lines
                    )
                    break
                except UnicodeDecodeError as e:
                    last_error = e
                    continue
                except Exception as e:
                    last_error = e
                    break
            
            if df is None:
                if last_error:
                    raise RuntimeError(f"Error reading CSV file: {str(last_error)}")
                else:
                    raise RuntimeError("Unable to read CSV file: could not determine encoding")
            
            # For CSV, we have a single "sheet" (the file itself)
            sheets = {"data": df}
            
        else:
            # Excel files (.xlsx, .xls, .xlsm)
            # Check for appropriate engine
            if suffix_lower in [".xlsx", ".xlsm"]:
                if openpyxl is None:
                    raise RuntimeError(
                        "openpyxl is required for .xlsx/.xlsm files. "
                        "Please install dependencies from requirements.txt"
                    )
                engine = "openpyxl"
            elif suffix_lower == ".xls":
                if xlrd is None:
                    raise RuntimeError(
                        "xlrd is required for .xls files. "
                        "Please install dependencies from requirements.txt"
                    )
                engine = "xlrd"
            else:
                raise ValueError(f"Unsupported spreadsheet format: {suffix_lower}")
            
            # Read all sheets
            sheets = pd.read_excel(
                str(input_path),
                sheet_name=None,  # Read all sheets
                engine=engine,
                dtype=str,  # Read all values as strings to preserve formatting
                na_filter=False  # Don't convert empty cells to NaN
            )
            
    except pd.errors.EmptyDataError:
        raise ValueError("Spreadsheet file is empty or has no readable data")
    except (RuntimeError, ValueError):
        # Re-raise our own errors
        raise
    except Exception as e:
        error_msg = str(e)
        raise RuntimeError(f"Error reading spreadsheet file: {error_msg}")
    
    if not sheets:
        raise ValueError("Spreadsheet file contains no data")
    
    # Check if all sheets are empty
    all_empty = True
    for sheet_name, df in sheets.items():
        if df is not None and not df.empty:
            all_empty = False
            break
    
    if all_empty:
        raise ValueError("Spreadsheet file contains no readable data")
    
    # Build markdown
    lines = [f"# {title}", ""]
    
    # Determine if this is a CSV file (no sheet headers needed)
    is_csv = suffix_lower == ".csv"
    
    # Convert each sheet to markdown table
    sheet_count = len(sheets)
    for sheet_name, df in sheets.items():
        if df is None or df.empty:
            continue
        
        # Add sheet header for Excel files with multiple sheets
        # Skip headers for CSV files (single data source, no sheet concept)
        if not is_csv and sheet_count > 1:
            table_md = _dataframe_to_markdown_table(df, sheet_name)
        else:
            table_md = _dataframe_to_markdown_table(df, None)
        
        if table_md:
            lines.append(table_md)
            lines.append("")  # Add blank line between sheets
    
    md = "\n".join(lines).strip() + "\n"
    return md


def file_to_markdown(file_path: str | Path) -> str:
    """
    Convert a file (PDF, DOCX, TXT, PPTX, CSV, or Excel) to Markdown format string.
    If the input is already Markdown (.md), return its content unchanged.
    
    Args:
        file_path: Path to the input file (PDF, DOCX, TXT, PPTX, CSV, Excel, or MD)
    
    Returns:
        Markdown content as a string
    
    Raises:
        FileNotFoundError: If the input file does not exist.
        ValueError: If the file type is not supported or file is empty/corrupted.
        RuntimeError: If required dependencies are missing.
    """
    input_path = Path(file_path).resolve()
    
    if not input_path.exists():
        raise FileNotFoundError(f"File not found: {input_path}")
    
    if not input_path.is_file():
        raise ValueError(f"Path is not a file: {input_path}")
    
    # Determine file type from extension
    suffix_lower = input_path.suffix.lower()
    
    # Handle markdown passthrough
    if suffix_lower == ".md":
        return input_path.read_text(encoding="utf-8")

    # Convert based on file type
    if suffix_lower == ".pdf":
        return convert_pdf_to_markdown(input_path)
    elif suffix_lower in [".docx", ".doc"]:
        return convert_docx_to_markdown(input_path)
    elif suffix_lower == ".txt":
        return convert_txt_to_markdown(input_path)
    elif suffix_lower == ".pptx":
        return convert_pptx_to_markdown(input_path)
    elif suffix_lower in [".csv", ".xlsx", ".xls", ".xlsm"]:
        return convert_spreadsheet_to_markdown(input_path)
    else:
        raise ValueError(
            f"Unsupported file type: {suffix_lower}. "
            f"Supported types: .pdf, .docx, .doc, .pptx, .txt, .md, .csv, .xlsx, .xls, .xlsm"
        )
